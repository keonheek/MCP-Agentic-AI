import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
NAVY = RGBColor(0x1F, 0x38, 0x64)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# Slide dimensions: widescreen 13.333 x 7.5 inches
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# --- Layout helpers ---

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, text: str, left, top, width, height,
                  font_name="Calibri", font_size=14, bold=False,
                  color=DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True) -> None:
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title_bar(slide, title: str):
    """Navy bar at top with white title text."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), SLIDE_W, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _add_body_text(slide, text: str, top_offset=Inches(1.3),
                   left=Inches(0.5), width=Inches(12.3), height=Inches(5.5)):
    txBox = slide.shapes.add_textbox(left, top_offset, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(14)
    run.font.color.rgb = DARK_GRAY


def _add_bullet_list(slide, items: list[str], top_offset=Inches(1.3),
                     left=Inches(0.5), width=Inches(12.3)):
    txBox = slide.shapes.add_textbox(left, top_offset, width, Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.name = "Calibri"
        run.font.size = Pt(13)
        run.font.color.rgb = DARK_GRAY


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars - 3] + "..."


# --- Individual slide builders ---

def _slide_01_title(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, WHITE)

    # Full-width navy background shape (top 60%)
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(4.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Company title
    company = state.get("company_description", "Company")[:80]
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = company
    run.font.name = "Calibri"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Problem subtitle
    problem = state.get("problem_statement", "")[:200]
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = problem
    run2.font.name = "Calibri"
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0xCC, 0xD6, 0xE8)

    # Footer
    txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.6))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "SME Business Diagnostic | Confidential"
    run3.font.name = "Calibri"
    run3.font.size = Pt(11)
    run3.font.color.rgb = DARK_GRAY


def _slide_02_exec_summary(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Executive Summary")

    final_recs = state.get("final_recommendations", [])
    benchmark_results = state.get("benchmark_results", {})

    top3 = final_recs[:3]
    key_finding = ""
    if benchmark_results:
        first_key = next(iter(benchmark_results))
        key_finding = _truncate(benchmark_results[first_key], 200)

    lines = ["Key Recommendations:"]
    for i, rec in enumerate(top3, 1):
        lines.append(f"  {i}. {rec.get('title', '')} — Impact: {rec.get('impact', 'medium').upper()}")

    if key_finding:
        lines.append("")
        lines.append("Key Market Finding:")
        lines.append(f"  {key_finding}")

    _add_bullet_list(slide, lines)


def _slide_03_problem_decomp(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Problem Decomposition")

    driver_tree = state.get("driver_tree", {})
    root = driver_tree.get("root", "Root problem")
    branches = driver_tree.get("branches", [])

    lines = [f"ROOT: {root}", ""]
    for branch in branches:
        lines.append(f"  [{branch.get('name', '')}]")
        for sub in branch.get("sub_branches", []):
            lines.append(f"      - {sub}")
        lines.append("")

    _add_bullet_list(slide, lines)


def _slide_04_hypotheses(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Hypotheses Investigated")

    hypotheses = state.get("hypotheses", [])
    lines = []
    for h in hypotheses:
        lines.append(f"  \u2610  {h}")
        lines.append("")

    _add_bullet_list(slide, lines if lines else ["No hypotheses recorded."])


def _slide_05_market_intel(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Market Intelligence")

    benchmark_results = state.get("benchmark_results", {})
    keys = list(benchmark_results.keys())[:3]

    lines = []
    for key in keys:
        text = _truncate(benchmark_results[key], 300)
        lines.append(f"[{key}]")
        lines.append(f"  {text}")
        lines.append("")

    _add_bullet_list(slide, lines if lines else ["No benchmark data available."])


def _slide_06_competitive(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Competitive Landscape")

    benchmark_results = state.get("benchmark_results", {})
    keys = list(benchmark_results.keys())[3:]

    lines = []
    for key in keys:
        text = _truncate(benchmark_results[key], 300)
        lines.append(f"[{key}]")
        lines.append(f"  {text}")
        lines.append("")

    if not lines:
        # Fall back to showing all if fewer than 4 benchmarks total
        all_keys = list(benchmark_results.keys())
        for key in all_keys:
            text = _truncate(benchmark_results[key], 200)
            lines.append(f"[{key}]: {text}")
            lines.append("")

    _add_bullet_list(slide, lines if lines else ["See previous slide for full competitive data."])


def _slide_07_gap_analysis(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Gap Analysis")

    benchmark_results = state.get("benchmark_results", {})
    driver_tree = state.get("driver_tree", {})
    branches = driver_tree.get("branches", [])

    lines = ["Area                        | Your Position  | Industry Benchmark",
             "-" * 65]

    for branch in branches:
        area = branch.get("name", "Unknown")[:25]
        bench_text = benchmark_results.get(branch.get("name", ""), "")
        # Extract a benchmark figure if present (heuristic: first number + % or ratio)
        import re
        nums = re.findall(r'\d+\.?\d*\s*%', bench_text)
        benchmark_fig = nums[0] if nums else "See benchmark data"
        lines.append(f"  {area:<27} | Under review   | {benchmark_fig}")

    _add_bullet_list(slide, lines)


def _slide_08_root_cause(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Root Cause Analysis")

    problem_type = state.get("problem_type", "unknown")
    hypotheses = state.get("hypotheses", [])
    driver_tree = state.get("driver_tree", {})
    benchmark_results = state.get("benchmark_results", {})

    top_hypothesis = hypotheses[0] if hypotheses else "No hypothesis identified."
    root = driver_tree.get("root", "")

    # Supporting evidence from first benchmark entry
    supporting = ""
    if benchmark_results:
        first_val = next(iter(benchmark_results.values()))
        supporting = _truncate(first_val, 250)

    lines = [
        f"Problem Classification: {problem_type.upper()}",
        "",
        f"Core Issue: {root}",
        "",
        "Primary Hypothesis:",
        f"  {top_hypothesis}",
        "",
        "Supporting Evidence:",
        f"  {supporting}",
    ]

    _add_bullet_list(slide, lines)


def _slide_09_recommendations(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Recommendations")

    final_recs = state.get("final_recommendations", [])

    lines = []
    for i, rec in enumerate(final_recs, 1):
        title = rec.get("title", f"Recommendation {i}")
        desc = _truncate(rec.get("description", ""), 180)
        impact = rec.get("impact", "medium").upper()
        feasibility = rec.get("feasibility", "medium").upper()
        lines.append(f"{i}. {title}")
        lines.append(f"   {desc}")
        lines.append(f"   Impact: {impact}  |  Feasibility: {feasibility}")
        lines.append("")

    _add_bullet_list(slide, lines if lines else ["No recommendations generated."])


def _slide_10_roadmap(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Implementation Roadmap")

    final_recs = state.get("final_recommendations", [])

    h1, h2, h3 = [], [], []
    for rec in final_recs:
        feasibility = rec.get("feasibility", "medium").lower()
        title = rec.get("title", "Untitled")
        if feasibility == "high":
            h1.append(title)
        elif feasibility == "medium":
            h2.append(title)
        else:
            h3.append(title)

    lines = [
        "H1 (0-3 months) — Quick Wins [High Feasibility]:",
    ] + [f"  - {t}" for t in h1] + (["  - (none)"] if not h1 else []) + [
        "",
        "H2 (3-6 months) — Core Initiatives [Medium Feasibility]:",
    ] + [f"  - {t}" for t in h2] + (["  - (none)"] if not h2 else []) + [
        "",
        "H3 (6-12 months) — Strategic Bets [Lower Feasibility]:",
    ] + [f"  - {t}" for t in h3] + (["  - (none)"] if not h3 else [])

    _add_bullet_list(slide, lines)


def _slide_11_roi(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "ROI & Success Metrics")

    final_recs = state.get("final_recommendations", [])
    top3 = final_recs[:3]

    # Heuristic KPI mapping based on impact/description keywords
    kpi_map = {
        "cost": "Cost reduction (%)",
        "revenue": "Revenue growth (%)",
        "efficiency": "OEE improvement (%)",
        "price": "ASP increase (%)",
        "customer": "Customer retention rate (%)",
        "market": "Market share gain (%)",
        "margin": "Operating margin improvement (pp)",
        "digital": "Process automation rate (%)",
    }

    lines = ["Recommendation                    | KPI Target               | Measurement",
             "-" * 80]
    for rec in top3:
        title = _truncate(rec.get("title", ""), 30)
        desc_lower = rec.get("description", "").lower()
        kpi = "Operating KPI improvement (%)"
        for keyword, kpi_label in kpi_map.items():
            if keyword in desc_lower or keyword in rec.get("title", "").lower():
                kpi = kpi_label
                break
        lines.append(f"  {title:<32} | {kpi:<24} | Monthly review")

    _add_bullet_list(slide, lines)


def _slide_12_next_steps(prs: Presentation, state: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _add_title_bar(slide, "Next Steps (30 Days)")

    final_recs = state.get("final_recommendations", [])
    top_rec = final_recs[0] if final_recs else {}

    title = top_rec.get("title", "Lead recommendation")
    desc = top_rec.get("description", "")

    # Derive 5 concrete 30-day actions from the top recommendation
    actions = [
        f"1. Assign workstream owner for '{title}'",
        f"2. Conduct internal data audit to baseline current performance",
        f"3. Identify 2-3 pilot units / departments to run proof-of-concept",
        f"4. Set 30-day KPI targets and establish tracking mechanism",
        f"5. Schedule weekly check-in — report progress to leadership by Day 30",
    ]

    lines = [f"Focused on: {title}", f"Context: {_truncate(desc, 150)}", ""] + actions

    _add_bullet_list(slide, lines)


# --- Main entry point ---

def generate_deck(state: dict, output_path: str = "output/consulting_deck.pptx") -> str:
    """Generate 12-slide pptx. Returns output_path."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_01_title(prs, state)
    _slide_02_exec_summary(prs, state)
    _slide_03_problem_decomp(prs, state)
    _slide_04_hypotheses(prs, state)
    _slide_05_market_intel(prs, state)
    _slide_06_competitive(prs, state)
    _slide_07_gap_analysis(prs, state)
    _slide_08_root_cause(prs, state)
    _slide_09_recommendations(prs, state)
    _slide_10_roadmap(prs, state)
    _slide_11_roi(prs, state)
    _slide_12_next_steps(prs, state)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"[deck_generator] Saved: {output_path}")
    return output_path


if __name__ == "__main__":
    # Fully offline test — no API keys needed
    sample_state = {
        "company_description": "Korean manufacturing SME, 400 employees, produces automotive components",
        "problem_statement": "영업이익률이 3년 연속 하락하고 있음. 원인이 뭔지, 어떻게 회복할 수 있는지 모름.",
        "problem_type": "operations",
        "driver_tree": {
            "root": "Declining operating margin",
            "branches": [
                {
                    "name": "Cost of Goods Sold",
                    "sub_branches": ["Raw material cost", "Labor cost"],
                },
                {
                    "name": "Revenue per Unit",
                    "sub_branches": ["Pricing power", "Customer mix"],
                },
                {
                    "name": "Operating Efficiency",
                    "sub_branches": ["Yield rate", "Overhead absorption"],
                },
            ],
        },
        "hypotheses": [
            "Hypothesis: Raw material cost increases are the primary driver of margin compression.",
            "Hypothesis: Pricing power has eroded due to competition from Chinese manufacturers.",
            "Hypothesis: Operational inefficiencies have worsened due to aging equipment.",
        ],
        "benchmark_results": {
            "Cost of Goods Sold": (
                "Korean auto parts manufacturers average COGS ratio of 72-75%. "
                "Raw material costs increased 18% YoY in 2023. Labor costs rose 5.1%."
            ),
            "Revenue per Unit": (
                "Average selling price for mid-tier auto components declined 3.2%. "
                "Hyundai/Kia pushed suppliers on cost reduction. Tier-2 suppliers face stronger pressure."
            ),
            "Operating Efficiency": (
                "Industry average OEE for Korean auto parts is 68%. Top quartile achieves 82%+. "
                "Smart factory adoption at 34% of mid-size manufacturers."
            ),
            "Market Outlook": (
                "EV transition accelerating — ICE component demand projected to decline 12% by 2027. "
                "Suppliers need to diversify into EV parts or face structural revenue decline."
            ),
        },
        "final_recommendations": [
            {
                "title": "Raw Material Hedging Program",
                "description": (
                    "Implement a 6-month forward contract hedging strategy for steel and aluminum. "
                    "Partner with a Korean commodity broker to lock in prices. "
                    "This directly addresses the 18% raw material cost spike."
                ),
                "impact": "high",
                "feasibility": "high",
            },
            {
                "title": "OEE Improvement Initiative",
                "description": (
                    "Deploy IoT sensors on top 5 production lines to track OEE in real-time. "
                    "Target improvement from estimated 62% to 72% within 6 months. "
                    "ROI payback estimated at 14 months based on industry benchmarks."
                ),
                "impact": "high",
                "feasibility": "medium",
            },
            {
                "title": "EV Component Diversification",
                "description": (
                    "Identify 2 EV component categories where existing tooling can be repurposed. "
                    "Apply for Korean government EV supply chain grant (deadline Q3). "
                    "Begin pilot production run for Tier-1 EV supplier qualification."
                ),
                "impact": "high",
                "feasibility": "medium",
            },
            {
                "title": "Pricing Renegotiation with Tier-1 Customers",
                "description": (
                    "Use raw material cost data to renegotiate pricing with Hyundai/Kia. "
                    "Present cost pass-through model with 6-month lag. "
                    "Industry peers have secured 2-4% ASP increases using this approach."
                ),
                "impact": "medium",
                "feasibility": "medium",
            },
            {
                "title": "Fixed Cost Reduction via Overhead Audit",
                "description": (
                    "Conduct a 30-day overhead absorption audit. "
                    "Identify underutilized shifts and consolidate production lines. "
                    "Target 8% SG&A reduction within 90 days."
                ),
                "impact": "medium",
                "feasibility": "high",
            },
        ],
    }

    out = generate_deck(sample_state, output_path="output/test_deck.pptx")
    print(f"Test deck saved to: {out}")
